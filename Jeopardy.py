import requests
from bs4 import BeautifulSoup
import random


C1, C2, C3, C4, C5, C6, C7, C8, C9, C10, C11, C12, F =  \
    {}, {}, {}, {}, {}, {}, {}, {}, {}, {}, {}, {}, {}

G = {1:C1, 2:C2, 3:C3, 4:C4, 5:C5, 6:C6, 7:C7, 8:C8, 9:C9, 10:C10, 11:C11, 12:C12, "f":F}


print()
op = int(input("Would you like to choose categories from a list (1), or decide randomly (2)? "))

if op == 1:
    print()
    C = []
    Qs = []
    As = []
    urls = []
    while len(C) < 50:
        url = "https://j-archive.com/showgame.php?game_id=" + str(random.randint(1,7400))
        urls.append(url)
        soup = BeautifulSoup(requests.get(url).text, "html.parser")
        clues = soup.find_all("td", {"class": "clue"})
        categories = list(map(lambda e: e.get_text(), soup.find_all("td", {"class": "category_name"})))
        for c in range(1,13):
            qs = []
            ans = []
            valid = True
            for q in range(1,6):
                d = soup.find(id="clue_" + c//7*"D" + "J_" + str(c%7 + c//7) + "_" + str(q))
                if d == None:
                    valid = False
                else:
                    qs.append(d.get_text())
                    ans.append(BeautifulSoup(clues[c//7*31+c%7+6*(q-1)-1].find(lambda e: e.name=="div"
                        and "onmouseover" in e.attrs)["onmouseover"], "html.parser").em.get_text())
            if valid:
                C.append(categories[c-1])
                Qs.append(qs)
                As.append(ans)

    for i in range(0,50):
        print("Option %s: " % str(i+1) + C[i])
    print()
    c = input("Enter 13 numbers corresponding to options above separated by spaces: ").split(" ")
    for i in range(1,13):
        G[i]["n"] = C[int(c[i-1])-1]
        for q in range(1,6):
            G[i][q] = Qs[int(c[i-1])-1][q-1]
            G[i]["a"+str(q)] = As[int(c[i-1])-1][q-1]
    G["f"]["n"] = C[int(c[12])-1]
    G["f"]["q"] = Qs[int(c[12])-1][2]
    G["f"]["a"] = As[int(c[12])-1][2]

if op == 2:
    C = []
    Qs = []
    As = []
    urls = []
    while len(C) < 13:
        url = "https://j-archive.com/showgame.php?game_id=" + str(random.randint(1,7400))
        urls.append(url)
        soup = BeautifulSoup(requests.get(url).text, "html.parser")
        clues = soup.find_all("td", {"class": "clue"})
        categories = list(map(lambda e: e.get_text(), soup.find_all("td", {"class": "category_name"})))
        for c in range(1,13):
            qs = []
            ans = []
            valid = True
            for q in range(1,6):
                d = soup.find(id="clue_" + c//7*"D" + "J_" + str(c%7 + c//7) + "_" + str(q))
                if d == None:
                    valid = False
                else:
                    qs.append(d.get_text())
                    ans.append(BeautifulSoup(clues[c//7*31+c%7+6*(q-1)-1].find(lambda e: e.name=="div"
                        and "onmouseover" in e.attrs)["onmouseover"], "html.parser").em.get_text())
            if valid:
                C.append(categories[c-1])
                Qs.append(qs)
                As.append(ans)

    for i in range(1,13):
        G[i]["n"] = C[i-1]
        for q in range(1,6):
            G[i][q] = Qs[i-1][q-1]
            G[i]["a"+str(q)] = As[i-1][q-1]
    G["f"]["n"] = C[12]
    G["f"]["q"] = Qs[12][2]
    G["f"]["a"] = As[12][2]


from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

prs = Presentation('Downloads/Jeopardy Template.pptx')

for i in range(1,13):
    print()
    print(G[i]["n"]+ ": ")

    slide = prs.slides[i+31*(i//7)]
    text_frame = slide.shapes[0].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = G[i]["n"]
    font = run.font
    font.name = 'Impact'
    font.size = Pt(96)
    font.bold = False
    font.color.rgb = RGBColor(255, 255, 255)

    slide = prs.slides[7+37*(i//7)]
    text_frame = slide.shapes[i-1-6*(i//7)].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = G[i]["n"]
    font = run.font
    font.name = 'Impact'
    font.size = Pt(20)
    font.bold = False
    font.color.rgb = RGBColor(255, 255, 255)

    for q in range(1,6):
        print("$" + str(200*q*(1+(i//7))) + ": " + G[i]["a"+str(q)].capitalize())
        slide = prs.slides[2+i*5+q+7*(i//7)]
        text_frame = slide.shapes[0].text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = G[i][q].upper()
        font = run.font
        font.name = 'Times New Roman'
        font.size = Pt(32)
        font.bold = True
        font.color.rgb = RGBColor(255, 255, 255)

print()
print("Final Jeopardy:")
print(G["f"]["a"])

slide = prs.slides[76]
text_frame = slide.shapes[0].text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = G["f"]["n"]
font = run.font
font.name = 'Impact'
font.size = Pt(96)
font.bold = False
font.color.rgb = RGBColor(255, 255, 255)

slide = prs.slides[77]
text_frame = slide.shapes[0].text_frame
text_frame.clear()
p = text_frame.paragraphs[0]
run = p.add_run()
run.text = G["f"]["q"].upper()
font = run.font
font.name = 'Times New Roman'
font.size = Pt(32)
font.bold = True
font.color.rgb = RGBColor(255, 255, 255)

print()
print("Daily Doubles:")
print("$" + str(random.randint(0,2)*200+600) + ": " + C[random.randint(0,5)])
print("$" + str(random.randint(0,2)*400+1200) + ": " + C[random.randint(6,8)])
print("$" + str(random.randint(0,2)*400+1200) + ": " + C[random.randint(9,11)])

print()
print("Source: ")
for url in urls:
    print(url)


prs.save('Downloads/Jeopardy Game.pptx')
